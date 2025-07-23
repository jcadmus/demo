import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import datetime

def abbreviate_amount(n):
    if n is None or pd.isna(n):
        return "$0"
    n = float(n)
    if abs(n) >= 1_000_000_000:
        return f"${n/1_000_000_000:.2f}B"
    elif abs(n) >= 1_000_000:
        return f"${n/1_000_000:.2f}M"
    elif abs(n) >= 1_000:
        return f"${n/1_000:.1f}K"
    else:
        return f"${n:.0f}"

def ppt_add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def ppt_add_chart_slide(prs, title, chart_img, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    pic = slide.shapes.add_picture(chart_img, Inches(1), Inches(1.3), Inches(8), Inches(4.5))
    if notes:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = notes
        tf.paragraphs[0].font.size = Pt(14)

def insight_card(msg, color="info"):
    c = {"info":"#232945", "warn":"#F39C12", "danger":"#E74C3C", "good":"#16A085"}
    st.markdown(f"<div style='background:{c.get(color, '#232945')};border-radius:12px;padding:14px 18px;color:white;font-size:1.13em;margin-bottom:12px'>{msg}</div>", unsafe_allow_html=True)

st.set_page_config(page_title="USAF Contract Command Center", layout="wide")

kpi_style = """
<style>
[data-testid="metric-container"] {
    background: #232945 !important;
    border-radius: 18px;
    padding: 30px 10px 30px 20px;
    margin: 10px 0 20px 0;
    color: #1ef1f7 !important;
    box-shadow: 0 4px 24px 0 rgba(30, 241, 247, 0.15);
    font-family: 'Segoe UI', sans-serif;
}
[data-testid="stMetricValue"] {
    font-size: 2.4rem;
    color: #1ef1f7 !important;
    font-weight: bold;
}
[data-testid="stMetricLabel"] {
    font-size: 1.1rem;
    color: #feee5b !important;
    font-weight: 600;
    opacity: 0.95;
}
</style>
"""
st.markdown(kpi_style, unsafe_allow_html=True)

st.markdown(
    "<h1 style='color:#1ef1f7; margin-bottom: 10px;'>USAF Contract Command Center</h1>",
    unsafe_allow_html=True
)
st.markdown("Upload a synthetic Air Force contract Excel file (`.xlsx`) to get started.")

uploaded_file = st.file_uploader(
    "Choose an Air Force Contract Excel file",
    type=["xlsx"],
    help="Upload your Air Force synthetic contracting data file."
)

if uploaded_file is not None:
    try:
        df = pd.read_excel(uploaded_file)
        date_cols = ["Effective Date", "Last Modified Date", "Ultimate Completion Date"]
        for col in date_cols:
            if col in df.columns:
                df[col] = pd.to_datetime(df[col])

        # Sidebar
        st.sidebar.header("🔎 Filter Data")
        centers = df['Contracting Center'].dropna().unique().tolist()
        vendors = df['Vendor/Recipient'].dropna().unique().tolist()
        award_types = df['Award Type'].dropna().unique().tolist()
        programs = df['Program'].dropna().unique().tolist()
        years = sorted(df['Effective Date'].dt.year.unique())
        selected_centers = st.sidebar.multiselect("Contracting Center", centers, default=centers)
        selected_award_types = st.sidebar.multiselect("Award Type", award_types, default=award_types)
        selected_programs = st.sidebar.multiselect("Program", programs, default=programs[:5])
        selected_years = st.sidebar.slider("Award Year Range", min_value=int(min(years)), max_value=int(max(years)), value=(int(min(years)), int(max(years))))
        selected_flag = st.sidebar.selectbox("Show Only Latest Mods?", ["All", "Only Latest"], index=0)
        search_vendor = st.sidebar.text_input("Search Vendor/Recipient")
        search_initiative = st.sidebar.text_input("Search Initiative")
        st.sidebar.markdown("---")
        st.sidebar.markdown("**Scenario: Project Closeout Rate (%)**")
        scenario_closeout_rate = st.sidebar.slider("Assume percent of contracts that will close out in next year:", 0, 100, 90)

        # Breadcrumb
        if "breadcrumbs" not in st.session_state:
            st.session_state.breadcrumbs = []
        drill_columns = ["Contracting Center", "Award Type", "Vendor/Recipient", "Program"]

        def show_breadcrumbs():
            if st.session_state.breadcrumbs:
                st.write("Drilldown Path:")
                cols = st.columns(len(st.session_state.breadcrumbs) + 1)
                for i, (col_name, col_value) in enumerate(st.session_state.breadcrumbs):
                    if cols[i].button(f"{col_name}: {col_value}", key=f"bc_{i}"):
                        st.session_state.breadcrumbs = st.session_state.breadcrumbs[:i+1]
                        st.experimental_rerun()
                if cols[-1].button("Home", key="bc_home"):
                    st.session_state.breadcrumbs = []
                    st.experimental_rerun()
        show_breadcrumbs()

        # Main filter logic
        filtered_df = df.copy()
        filtered_df = filtered_df[
            filtered_df['Contracting Center'].isin(selected_centers) &
            filtered_df['Award Type'].isin(selected_award_types) &
            filtered_df['Program'].isin(selected_programs) &
            (filtered_df['Effective Date'].dt.year.between(selected_years[0], selected_years[1]))
        ]
        if selected_flag == "Only Latest":
            filtered_df = filtered_df[filtered_df['Record Flag'] == "latest"]
        if search_vendor:
            filtered_df = filtered_df[filtered_df['Vendor/Recipient'].str.contains(search_vendor, case=False, na=False)]
        if search_initiative:
            filtered_df = filtered_df[filtered_df['Initiative'].str.contains(search_initiative, case=False, na=False)]
        for col, val in st.session_state.breadcrumbs:
            filtered_df = filtered_df[filtered_df[col] == val]

        # ------ FIX: Ensure FY column is always present after filtering ------
        filtered_df['FY'] = filtered_df['Effective Date'].apply(
            lambda d: d.year + 1 if d.month >= 10 else d.year
        )

        drill_candidates = [c for c in drill_columns if c not in [col for col, val in st.session_state.breadcrumbs]]
        if drill_candidates and not filtered_df.empty:
            drill_col = st.selectbox("Drilldown: Pick a category to filter deeper", [""] + drill_candidates)
            if drill_col:
                drill_values = filtered_df[drill_col].dropna().unique().tolist()
                drill_value = st.selectbox(f"Select a {drill_col}", [""] + drill_values)
                if drill_value and (drill_col, drill_value) not in st.session_state.breadcrumbs:
                    st.session_state.breadcrumbs.append((drill_col, drill_value))
                    st.experimental_rerun()

        # ---------- KPIs ----------
        kpi1, kpi2, kpi3, kpi4, kpi5 = st.columns(5)
        kpi1.metric("Total Obligated Amount", abbreviate_amount(filtered_df['Action Obligated Amount'].sum()))
        kpi2.metric("Total Contracts/Orders", filtered_df['Contract Number'].nunique())
        kpi3.metric("Total Contract Value", abbreviate_amount(filtered_df['Total Actions & Options (Total Contract Value)'].sum()))
        kpi4.metric("Unliquidated Amount", abbreviate_amount(filtered_df['Unliquidated Amount'].sum()))
        fy25_start = pd.Timestamp("2024-10-01")
        fy25_end = pd.Timestamp("2025-09-30")
        fy25_expiring = filtered_df[
            (filtered_df['Ultimate Completion Date'] >= fy25_start) &
            (filtered_df['Ultimate Completion Date'] <= fy25_end)
        ]['Contract Number'].nunique()
        kpi5.metric("Contracts Expiring in FY25", fy25_expiring)
        st.markdown("---")

        # ---------- Automated Insight Cards ----------
        try:
            top_prog = filtered_df.groupby('Program')['Action Obligated Amount'].sum().idxmax()
            top_prog_amt = abbreviate_amount(filtered_df.groupby('Program')['Action Obligated Amount'].sum().max())
            insight_card(f"**Top Program:** {top_prog} ({top_prog_amt})", "good")
        except:
            pass
        try:
            top_vendor = filtered_df.groupby('Vendor/Recipient')['Action Obligated Amount'].sum().idxmax()
            top_vendor_amt = abbreviate_amount(filtered_df.groupby('Vendor/Recipient')['Action Obligated Amount'].sum().max())
            insight_card(f"**Largest Vendor:** {top_vendor} ({top_vendor_amt})", "good")
        except:
            pass
        rising_unliq = filtered_df.groupby('FY')['Unliquidated Amount'].sum().pct_change().fillna(0).iloc[-1]
        if rising_unliq > 0.10:
            insight_card(f"**Warning:** Unliquidated obligations are rising this year ({rising_unliq:.1%})!", "warn")
        exp_spike = filtered_df['Ultimate Completion Date'].dt.year.value_counts().sort_index().diff().max()
        if exp_spike and exp_spike > 50:
            insight_card(f"**Upcoming Expiry Spike:** Contracts expiring spike by {int(exp_spike)} in upcoming FY!", "warn")

        # ---------- Obligated by Center ----------
        st.subheader("Obligated Amount by Contracting Center")
        center_ob = filtered_df.groupby('Contracting Center')['Action Obligated Amount'].sum().reset_index()
        fig_bar = px.bar(
            center_ob, x='Contracting Center', y='Action Obligated Amount',
            color='Contracting Center', title="Obligated by Center"
        )
        fig_bar.update_traces(
            text=center_ob['Action Obligated Amount'].apply(abbreviate_amount),
            textposition='outside',
            hovertemplate='%{x}: %{text}<extra></extra>'
        )
        st.plotly_chart(fig_bar, use_container_width=True)
        st.markdown("---")

        # ---------- Funding Forecast/History by Fiscal Year ----------
        st.subheader("Obligated Amount by Fiscal Year (Forecast & History)")
        fy_ob = filtered_df.groupby('FY')['Action Obligated Amount'].sum().reset_index()
        fig_fy = px.bar(fy_ob, x='FY', y='Action Obligated Amount', title="Obligated by Fiscal Year")
        fig_fy.update_traces(
            text=fy_ob['Action Obligated Amount'].apply(abbreviate_amount), textposition='outside'
        )
        st.plotly_chart(fig_fy, use_container_width=True)

        # ---------- Expiring Contracts Timeline ----------
        st.subheader("Contracts Expiring per Quarter")
        filtered_df['Exp_Q'] = filtered_df['Ultimate Completion Date'].dt.to_period('Q').astype(str)
        expiring = filtered_df.groupby('Exp_Q')['Contract Number'].nunique().reset_index().rename(columns={"Contract Number": "Expiring Contracts"})
        fig_exp = px.bar(expiring, x='Exp_Q', y='Expiring Contracts', title="Contracts Expiring per Quarter")
        st.plotly_chart(fig_exp, use_container_width=True)

        # ---------- Contract Closeout Status (Simulated) ----------
        st.subheader("Contract Completion Closeout Status")
        today = pd.Timestamp.today()
        if "Closeout Status" not in filtered_df.columns:
            filtered_df['Closeout Status'] = np.where(
                filtered_df['Ultimate Completion Date'] < today - pd.Timedelta(days=90), 'Y',
                np.where(filtered_df['Ultimate Completion Date'] < today, 'N', 'Null')
            )
        projected = filtered_df.copy()
        projected['Closeout Status'] = np.where(
            (projected['Closeout Status'] == 'N') & (projected['Ultimate Completion Date'] < today + pd.Timedelta(days=365)),
            np.where(np.random.rand(len(projected)) < scenario_closeout_rate/100, 'Y', 'N'),
            projected['Closeout Status']
        )
        closeout_counts = projected.groupby(['FY', 'Closeout Status'])['Contract Number'].nunique().reset_index()
        all_statuses = ['Y', 'N', 'Null']
        count_pivot = closeout_counts.pivot(index='FY', columns='Closeout Status', values='Contract Number').fillna(0).reindex(columns=all_statuses, fill_value=0)
        count_pivot['Total'] = count_pivot.sum(axis=1)
        pct_display = count_pivot[all_statuses].div(count_pivot['Total'], axis=0).reset_index()*100
        # 100% stacked bar (percent)
        st.markdown("**Percent Closeout by FY**")
        fig_close_pct = px.bar(
            pct_display, x="FY", y=all_statuses,
            title="Contract Completion Closeout %",
            labels={"value": "% of Contracts", "variable": "Closeout Status"},
            barmode="stack"
        )
        for s in all_statuses:
            fig_close_pct.for_each_trace(lambda t: t.update(texttemplate="%{y:.1f}%", textposition='inside'))
        st.plotly_chart(fig_close_pct, use_container_width=True)
        # Count stacked
        st.markdown("**Number of Contracts Closed Out by FY**")
        count_pivot_display = count_pivot.reset_index()
        fig_close_count = px.bar(
            count_pivot_display, x="FY", y=all_statuses,
            title="Contract Completion Closeout Number of Contracts",
            labels={"value": "# of Contracts", "variable": "Closeout Status"},
            barmode="stack"
        )
        for s in all_statuses:
            fig_close_count.for_each_trace(lambda t: t.update(texttemplate="%{y:,.0f}", textposition='inside'))
        st.plotly_chart(fig_close_count, use_container_width=True)
        # Slope line for % closed
        if 'Y' in pct_display.columns:
            st.markdown("**Closeout % Over Time (Line)**")
            fig_slope = px.line(
                pct_display, x="FY", y='Y', markers=True,
                title="Percent of Contracts Closed Over Time"
            )
            fig_slope.update_traces(mode='markers+lines+text', text=pct_display["Y"].round(1).astype(str) + "%", textposition="top center")
            st.plotly_chart(fig_slope, use_container_width=True)

        # ---------- Top Programs/Initiatives ----------
        st.subheader("Top 10 Programs & Initiatives by Funding")
        prog_pareto = filtered_df.groupby('Program')['Action Obligated Amount'].sum().reset_index().sort_values('Action Obligated Amount', ascending=False).head(10)
        fig_prog_bar = px.bar(
            prog_pareto, x='Action Obligated Amount', y='Program', orientation='h', title="Top 10 Programs by Obligated"
        )
        fig_prog_bar.update_traces(
            text=prog_pareto['Action Obligated Amount'].apply(abbreviate_amount),
            textposition='outside',
            hovertemplate='%{y}: %{text}<extra></extra>'
        )
        st.plotly_chart(fig_prog_bar, use_container_width=True)

        # ---------- Unliquidated Aging ----------
        st.subheader("Unliquidated Amount Aging")
        filtered_df['Contract Age (Years)'] = ((pd.to_datetime('today') - filtered_df['Effective Date']).dt.days / 365).round(1)
        age_bins = pd.cut(filtered_df['Contract Age (Years)'], bins=[0,1,2,3,5,10,100], labels=['<1yr','1-2yr','2-3yr','3-5yr','5-10yr','10+yr'])
        aging = filtered_df.groupby(age_bins)['Unliquidated Amount'].sum().reset_index()
        fig_age = px.bar(aging, x='Contract Age (Years)', y='Unliquidated Amount', title="Unliquidated Amount by Contract Age Bin")
        fig_age.update_traces(text=aging['Unliquidated Amount'].apply(abbreviate_amount), textposition='outside')
        st.plotly_chart(fig_age, use_container_width=True)

        # ---------- Award Cycle Time ----------
        st.subheader("Contract Cycle Time (Days) by Contracting Center")
        filtered_df['Cycle Time (Days)'] = (filtered_df['Last Modified Date'] - filtered_df['Effective Date']).dt.days
        cycle_box = px.box(filtered_df, x='Contracting Center', y='Cycle Time (Days)', points='all', title="Cycle Time by Center")
        st.plotly_chart(cycle_box, use_container_width=True)

        # ---------- Vendor/Recipient Concentration ----------
        st.subheader("Vendor/Recipient Share of Obligated Amount")
        vendor_pareto = filtered_df.groupby('Vendor/Recipient')['Action Obligated Amount'].sum().reset_index().sort_values('Action Obligated Amount', ascending=False).head(10)
        fig_vendor_pareto = px.bar(
            vendor_pareto, x='Action Obligated Amount', y='Vendor/Recipient', orientation='h', title="Top 10 Vendors by Obligated"
        )
        fig_vendor_pareto.update_traces(
            text=vendor_pareto['Action Obligated Amount'].apply(abbreviate_amount),
            textposition='outside',
            hovertemplate='%{y}: %{text}<extra></extra>'
        )
        st.plotly_chart(fig_vendor_pareto, use_container_width=True)

        vendor_share = filtered_df.groupby('Vendor/Recipient')['Action Obligated Amount'].sum().reset_index()
        fig_vendor_pie = px.pie(
            vendor_share, names='Vendor/Recipient', values='Action Obligated Amount', title="Vendor/Recipient Funding Share"
        )
        st.plotly_chart(fig_vendor_pie, use_container_width=True)

        # ---------- Scatter: Obligated Amount vs. Contract Duration ----------
        st.subheader("Obligated Amount vs. Contract Duration")
        filtered_df['Duration (Months)'] = ((filtered_df['Ultimate Completion Date'] - filtered_df['Effective Date']).dt.days / 30.44).round(1)
        fig_scatter = px.scatter(
            filtered_df, x='Duration (Months)', y='Action Obligated Amount', color='Program',
            hover_data=['Contract Number', 'Vendor/Recipient'], title="Obligated vs. Duration (by Program)"
        )
        st.plotly_chart(fig_scatter, use_container_width=True)

        # ---------- Awards by Type/Program Over Time ----------
        st.subheader("Awards by Type and Program Over Time")
        awards_time = filtered_df.copy()
        awards_time['FY'] = awards_time['Effective Date'].apply(lambda d: d.year + 1 if d.month >= 10 else d.year)
        area = px.area(
            awards_time, x='FY', y='Action Obligated Amount', color='Award Type',
            line_group='Program', title="Obligated by Award Type (Stacked by FY)"
        )
        st.plotly_chart(area, use_container_width=True)

        # ---------- Program/Center Heatmap ----------
        st.subheader("Obligated vs. Unliquidated by Program and Center (Heatmap)")
        heatmap_df = filtered_df.groupby(['Program', 'Contracting Center']).agg(
            {'Action Obligated Amount':'sum','Unliquidated Amount':'sum'}).reset_index()
        heatmap = px.density_heatmap(
            heatmap_df, x='Program', y='Contracting Center',
            z='Unliquidated Amount', color_continuous_scale='Blues', 
            title="Unliquidated Amount by Program/Center"
        )
        st.plotly_chart(heatmap, use_container_width=True)

        # ---------- Drillable Data Table with Download ----------
        st.subheader("Filtered Data Table (Downloadable)")
        st.dataframe(filtered_df)
        csv = filtered_df.to_csv(index=False).encode('utf-8')
        st.download_button("Download Filtered Data as CSV", data=csv, file_name='filtered_usaf_contracts.csv', mime='text/csv')

        # ---------- Export to PowerPoint ----------
        st.subheader("Export Key Metrics & Charts to PowerPoint")
        if st.button("Export to PowerPoint"):
            prs = Presentation()
            ppt_add_title_slide(prs, "USAF Contract Dashboard", f"Report generated: {datetime.datetime.now().strftime('%b %d, %Y %H:%M')}")
            # Save and insert plots as images
            chart_figs = [
                ("Obligated by Center", fig_bar),
                ("Obligated by Fiscal Year", fig_fy),
                ("Contracts Expiring per Quarter", fig_exp),
                ("Closeout % by Year", fig_close_pct),
                ("Top Programs by Obligated", fig_prog_bar)
            ]
            for title, fig in chart_figs:
                img_bytes = BytesIO()
                fig.write_image(img_bytes, format="png")
                img_bytes.seek(0)
                ppt_add_chart_slide(prs, title, img_bytes)
            # KPIs Slide
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Key Metrics"
            text = f"""
Total Obligated: {abbreviate_amount(filtered_df['Action Obligated Amount'].sum())}
Total Contracts: {filtered_df['Contract Number'].nunique()}
Total Value: {abbreviate_amount(filtered_df['Total Actions & Options (Total Contract Value)'].sum())}
Unliquidated: {abbreviate_amount(filtered_df['Unliquidated Amount'].sum())}
Contracts Expiring in FY25: {fy25_expiring}
"""
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(2))
            tf = txBox.text_frame
            tf.text = text
            tf.paragraphs[0].font.size = Pt(20)
            ppt_bytes = BytesIO()
            prs.save(ppt_bytes)
            ppt_bytes.seek(0)
            st.download_button(
                label="Download PowerPoint",
                data=ppt_bytes,
                file_name="usaf_contract_dashboard.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        st.markdown("---")
        st.caption("Use all filters, scenario sliders, and drilldowns. Export your analytics to CSV or PPTX for executive briefings and program reviews.")

    except Exception as e:
        st.error(f"Error loading or processing file: {e}")

else:
    st.info("Awaiting file upload. Please upload an Excel file to continue.")
